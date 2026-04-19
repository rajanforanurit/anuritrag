"""
services/document_loader.py — Universal document text extractor.

Supported file types (industry-standard coverage):
  Documents    : PDF, DOCX, DOC, TXT, RTF, ODT
  Spreadsheets : XLSX, XLS, ODS, CSV, TSV
  Presentations: PPTX, PPT
  Web/Markup   : HTML, HTM, XML, MD, MARKDOWN, RST
  Data         : JSON, JSONL, YAML, YML, TOML
  Code         : PY, JS, TS, JSX, TSX, JAVA, CPP, C, H, CS, GO, RB, PHP, SWIFT, KT, R, SQL, SH, BASH, PS1
  eBook        : EPUB
  Email        : EML
"""

from __future__ import annotations

import json
import logging
from dataclasses import dataclass, field
from pathlib import Path
from typing import List, Optional

logger = logging.getLogger(__name__)

# ── All supported extensions ────────────────────────────────────────────────────

SUPPORTED_EXTENSIONS: set = {
    # Documents
    ".pdf", ".docx", ".doc", ".txt", ".rtf", ".odt",
    # Spreadsheets
    ".xlsx", ".xls", ".ods", ".csv", ".tsv",
    # Presentations
    ".pptx", ".ppt",
    # Web / Markup
    ".html", ".htm", ".xml", ".md", ".markdown", ".rst",
    # Data formats
    ".json", ".jsonl", ".yaml", ".yml", ".toml",
    # Code files
    ".py", ".js", ".ts", ".jsx", ".tsx",
    ".java", ".cpp", ".c", ".h", ".cs",
    ".go", ".rb", ".php", ".swift", ".kt",
    ".r", ".sql", ".sh", ".bash", ".ps1",
    # eBook
    ".epub",
    # Email
    ".eml",
}


# ── Raw document dataclass ──────────────────────────────────────────────────────

@dataclass
class RawDocument:
    file_path: Path
    doc_id: str
    pages: List[dict]           # list of {"page": int, "text": str, ...}
    total_pages: int = 0
    source_type: str = "local"
    extra_metadata: dict = field(default_factory=dict)

    @property
    def full_text(self) -> str:
        return "\n".join(
            p.get("text", "") for p in self.pages if p.get("text")
        )


# ── Per-format extractors ───────────────────────────────────────────────────────

def _load_pdf(path: Path) -> List[dict]:
    """PDF: try pypdf first, fall back to PyMuPDF for scanned/complex PDFs."""
    try:
        from pypdf import PdfReader
        reader = PdfReader(str(path))
        pages  = []
        for i, page in enumerate(reader.pages, start=1):
            text = (page.extract_text() or "").strip()
            if text:
                pages.append({"page": i, "text": text})
        if pages:
            return pages
    except Exception:
        pass

    try:
        import fitz  # PyMuPDF
        doc   = fitz.open(str(path))
        pages = []
        for i, page in enumerate(doc, start=1):
            text = page.get_text("text").strip()
            if text:
                pages.append({"page": i, "text": text})
        return pages
    except Exception as exc:
        logger.error("PDF extraction failed for %s: %s", path.name, exc)
        return []


def _load_docx(path: Path) -> List[dict]:
    from docx import Document
    doc   = Document(str(path))
    lines = [p.text.strip() for p in doc.paragraphs if p.text.strip()]
    for table in doc.tables:
        for row in table.rows:
            row_text = " | ".join(c.text.strip() for c in row.cells if c.text.strip())
            if row_text:
                lines.append(row_text)
    text = "\n".join(lines)
    return [{"page": 1, "text": text}] if text else []


def _load_doc(path: Path) -> List[dict]:
    """Legacy .doc — convert via LibreOffice."""
    try:
        import subprocess, tempfile
        with tempfile.TemporaryDirectory() as tmp:
            subprocess.run(
                ["libreoffice", "--headless", "--convert-to", "txt:Text",
                 str(path), "--outdir", tmp],
                capture_output=True, timeout=30,
            )
            txt_file = Path(tmp) / (path.stem + ".txt")
            if txt_file.exists():
                text = txt_file.read_text(encoding="utf-8", errors="replace").strip()
                return [{"page": 1, "text": text}] if text else []
    except Exception:
        pass
    try:
        text = path.read_text(encoding="latin-1", errors="replace").strip()
        return [{"page": 1, "text": text}] if text else []
    except Exception:
        return []


def _load_txt(path: Path) -> List[dict]:
    for enc in ("utf-8", "utf-16", "latin-1", "cp1252"):
        try:
            text = path.read_text(encoding=enc, errors="strict").strip()
            return [{"page": 1, "text": text}] if text else []
        except (UnicodeDecodeError, LookupError):
            continue
    text = path.read_text(encoding="utf-8", errors="replace").strip()
    return [{"page": 1, "text": text}] if text else []


def _load_rtf(path: Path) -> List[dict]:
    try:
        from striprtf.striprtf import rtf_to_text
        raw  = path.read_text(encoding="utf-8", errors="replace")
        text = rtf_to_text(raw).strip()
        return [{"page": 1, "text": text}] if text else []
    except ImportError:
        logger.warning("striprtf not installed — reading RTF as plain text")
        return _load_txt(path)


def _load_odt(path: Path) -> List[dict]:
    try:
        import zipfile
        from bs4 import BeautifulSoup
        with zipfile.ZipFile(str(path)) as zf:
            xml = zf.read("content.xml")
        soup = BeautifulSoup(xml, "xml")
        text = soup.get_text(separator="\n").strip()
        return [{"page": 1, "text": text}] if text else []
    except Exception as exc:
        logger.error("ODT extraction failed for %s: %s", path.name, exc)
        return []


def _load_pptx(path: Path) -> List[dict]:
    from pptx import Presentation
    prs   = Presentation(str(path))
    pages = []
    for i, slide in enumerate(prs.slides, start=1):
        lines = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                lines.append(shape.text.strip())
            if shape.has_table:
                for row in shape.table.rows:
                    row_text = " | ".join(
                        c.text.strip() for c in row.cells if c.text.strip()
                    )
                    if row_text:
                        lines.append(row_text)
        text = "\n".join(lines)
        if text:
            pages.append({"page": i, "text": text})
    return pages


def _load_ppt(path: Path) -> List[dict]:
    """Legacy .ppt — convert via LibreOffice."""
    try:
        import subprocess, tempfile
        with tempfile.TemporaryDirectory() as tmp:
            subprocess.run(
                ["libreoffice", "--headless", "--convert-to", "pptx",
                 str(path), "--outdir", tmp],
                capture_output=True, timeout=60,
            )
            pptx_file = Path(tmp) / (path.stem + ".pptx")
            if pptx_file.exists():
                return _load_pptx(pptx_file)
    except Exception:
        pass
    return []


def _load_excel(path: Path) -> List[dict]:
    import pandas as pd
    ext    = path.suffix.lower()
    engine = "odf" if ext == ".ods" else ("xlrd" if ext == ".xls" else "openpyxl")
    pages  = []
    try:
        xf = pd.ExcelFile(path, engine=engine)
        for i, sheet_name in enumerate(xf.sheet_names, start=1):
            df = pd.read_excel(path, sheet_name=sheet_name, engine=engine)
            if df.empty:
                continue
            header = " | ".join(str(c) for c in df.columns)
            rows   = df.fillna("").astype(str).apply(lambda r: " | ".join(r), axis=1)
            text   = header + "\n" + "\n".join(rows)
            pages.append({"page": i, "sheet_name": str(sheet_name), "text": text.strip()})
    except Exception as exc:
        logger.error("Excel extraction failed for %s: %s", path.name, exc)
    return pages


def _load_csv(path: Path) -> List[dict]:
    import pandas as pd
    try:
        df = None
        for sep in (",", ";", "\t", "|"):
            try:
                df = pd.read_csv(path, sep=sep, encoding="utf-8", errors="replace")
                if len(df.columns) > 1:
                    break
            except Exception:
                continue
        if df is None or df.empty:
            return []
        header = " | ".join(str(c) for c in df.columns)
        rows   = df.fillna("").astype(str).apply(lambda r: " | ".join(r), axis=1)
        text   = header + "\n" + "\n".join(rows)
        return [{"page": 1, "text": text.strip()}]
    except Exception as exc:
        logger.error("CSV extraction failed for %s: %s", path.name, exc)
        return []


def _load_tsv(path: Path) -> List[dict]:
    import pandas as pd
    try:
        df = pd.read_csv(path, sep="\t", encoding="utf-8", errors="replace")
        if df.empty:
            return []
        header = " | ".join(str(c) for c in df.columns)
        rows   = df.fillna("").astype(str).apply(lambda r: " | ".join(r), axis=1)
        text   = header + "\n" + "\n".join(rows)
        return [{"page": 1, "text": text.strip()}]
    except Exception as exc:
        logger.error("TSV extraction failed for %s: %s", path.name, exc)
        return []


def _load_json(path: Path) -> List[dict]:
    try:
        with open(path, "r", encoding="utf-8", errors="replace") as f:
            data = json.load(f)
        text = json.dumps(data, indent=2, ensure_ascii=False)
        return [{"page": 1, "text": text}] if text else []
    except json.JSONDecodeError:
        return _load_jsonl(path)
    except Exception as exc:
        logger.error("JSON extraction failed for %s: %s", path.name, exc)
        return []


def _load_jsonl(path: Path) -> List[dict]:
    lines = []
    try:
        with open(path, "r", encoding="utf-8", errors="replace") as f:
            for line in f:
                line = line.strip()
                if not line:
                    continue
                try:
                    obj = json.loads(line)
                    lines.append(json.dumps(obj, ensure_ascii=False))
                except json.JSONDecodeError:
                    lines.append(line)
        text = "\n".join(lines)
        return [{"page": 1, "text": text}] if text else []
    except Exception as exc:
        logger.error("JSONL extraction failed for %s: %s", path.name, exc)
        return []


def _load_yaml(path: Path) -> List[dict]:
    try:
        import yaml
        with open(path, "r", encoding="utf-8", errors="replace") as f:
            data = yaml.safe_load(f)
        text = yaml.dump(data, allow_unicode=True, default_flow_style=False)
        return [{"page": 1, "text": text.strip()}] if text else []
    except ImportError:
        return _load_txt(path)
    except Exception as exc:
        logger.error("YAML extraction failed for %s: %s", path.name, exc)
        return []


def _load_toml(path: Path) -> List[dict]:
    try:
        try:
            import tomllib                  # Python 3.11+
        except ImportError:
            import tomli as tomllib         # pip install tomli
        with open(path, "rb") as f:
            data = tomllib.load(f)
        text = json.dumps(data, indent=2, ensure_ascii=False)
        return [{"page": 1, "text": text}] if text else []
    except Exception:
        return _load_txt(path)


def _load_html(path: Path) -> List[dict]:
    import re
    from bs4 import BeautifulSoup
    html = path.read_text(encoding="utf-8", errors="replace")
    soup = BeautifulSoup(html, "html.parser")
    for tag in soup(["script", "style", "noscript", "head"]):
        tag.decompose()
    text = soup.get_text(separator="\n").strip()
    text = re.sub(r"\n{3,}", "\n\n", text)
    return [{"page": 1, "text": text}] if text else []


def _load_xml(path: Path) -> List[dict]:
    from bs4 import BeautifulSoup
    xml  = path.read_text(encoding="utf-8", errors="replace")
    soup = BeautifulSoup(xml, "xml")
    text = soup.get_text(separator="\n").strip()
    return [{"page": 1, "text": text}] if text else []


def _load_markdown(path: Path) -> List[dict]:
    text = path.read_text(encoding="utf-8", errors="replace").strip()
    return [{"page": 1, "text": text}] if text else []


def _load_rst(path: Path) -> List[dict]:
    return _load_txt(path)


def _load_code(path: Path) -> List[dict]:
    """Source code files — read as plain text, tag with language."""
    text = path.read_text(encoding="utf-8", errors="replace").strip()
    if text:
        lang = path.suffix.lstrip(".")
        return [{"page": 1, "text": f"# Language: {lang}\n\n{text}"}]
    return []


def _load_epub(path: Path) -> List[dict]:
    try:
        import zipfile
        from bs4 import BeautifulSoup
        pages = []
        with zipfile.ZipFile(str(path)) as zf:
            html_files = sorted(
                n for n in zf.namelist()
                if n.endswith((".html", ".xhtml", ".htm"))
                and not n.startswith("__MACOSX")
            )
            for i, name in enumerate(html_files, start=1):
                html = zf.read(name).decode("utf-8", errors="replace")
                soup = BeautifulSoup(html, "html.parser")
                for tag in soup(["script", "style"]):
                    tag.decompose()
                text = soup.get_text(separator="\n").strip()
                if text:
                    pages.append({"page": i, "text": text})
        return pages
    except Exception as exc:
        logger.error("EPUB extraction failed for %s: %s", path.name, exc)
        return []


def _load_eml(path: Path) -> List[dict]:
    """Email files (.eml) — extract headers + body text."""
    import email
    from email import policy
    try:
        raw = path.read_bytes()
        msg = email.message_from_bytes(raw, policy=policy.default)
        parts = []
        for hdr in ("From", "To", "Subject", "Date"):
            val = msg.get(hdr, "")
            if val:
                parts.append(f"{hdr}: {val}")
        parts.append("")
        if msg.is_multipart():
            for part in msg.walk():
                ct = part.get_content_type()
                if ct == "text/plain":
                    parts.append(part.get_content())
                elif ct == "text/html":
                    from bs4 import BeautifulSoup
                    soup = BeautifulSoup(part.get_content(), "html.parser")
                    parts.append(soup.get_text(separator="\n"))
        else:
            parts.append(msg.get_content())
        text = "\n".join(str(p) for p in parts).strip()
        return [{"page": 1, "text": text}] if text else []
    except Exception as exc:
        logger.error("EML extraction failed for %s: %s", path.name, exc)
        return []


# ── Extractor dispatch map ──────────────────────────────────────────────────────

_EXTRACTOR_MAP = {
    # Documents
    ".pdf":      _load_pdf,
    ".docx":     _load_docx,
    ".doc":      _load_doc,
    ".txt":      _load_txt,
    ".rtf":      _load_rtf,
    ".odt":      _load_odt,
    # Spreadsheets
    ".xlsx":     _load_excel,
    ".xls":      _load_excel,
    ".ods":      _load_excel,
    ".csv":      _load_csv,
    ".tsv":      _load_tsv,
    # Presentations
    ".pptx":     _load_pptx,
    ".ppt":      _load_ppt,
    # Web / Markup
    ".html":     _load_html,
    ".htm":      _load_html,
    ".xml":      _load_xml,
    ".md":       _load_markdown,
    ".markdown": _load_markdown,
    ".rst":      _load_rst,
    # Data formats
    ".json":     _load_json,
    ".jsonl":    _load_jsonl,
    ".yaml":     _load_yaml,
    ".yml":      _load_yaml,
    ".toml":     _load_toml,
    # Code files
    ".py":    _load_code, ".js":    _load_code, ".ts":    _load_code,
    ".jsx":   _load_code, ".tsx":   _load_code, ".java":  _load_code,
    ".cpp":   _load_code, ".c":     _load_code, ".h":     _load_code,
    ".cs":    _load_code, ".go":    _load_code, ".rb":    _load_code,
    ".php":   _load_code, ".swift": _load_code, ".kt":    _load_code,
    ".r":     _load_code, ".sql":   _load_code, ".sh":    _load_code,
    ".bash":  _load_code, ".ps1":   _load_code,
    # eBook
    ".epub":  _load_epub,
    # Email
    ".eml":   _load_eml,
}


# ── DocumentLoader class ────────────────────────────────────────────────────────

class DocumentLoader:

    def load_from_directory(
        self,
        root: Path,
        source_type: str = "local",
        extra_metadata: Optional[dict] = None,
    ) -> List[RawDocument]:
        if not root.exists() or not root.is_dir():
            raise ValueError(f"Invalid directory path: {root}")

        files = [
            f for f in root.rglob("*")
            if f.is_file() and f.suffix.lower() in SUPPORTED_EXTENSIONS
        ]

        if not files:
            logger.warning("No supported files found in: %s", root)
            return []

        results = []
        for file_path in sorted(files):
            doc = self._extract(file_path, source_type=source_type, extra_metadata=extra_metadata)
            if doc:
                results.append(doc)
        return results

    def scan_directory(self, root: Path) -> List[dict]:
        """
        Lightweight scan — returns file metadata without extracting text.
        Used by /scan-directory and /ingest/local-directory (for file count).
        Returns list of dicts: File, Type, Size (KB), Subfolder.
        """
        if not root.exists() or not root.is_dir():
            return []

        rows: List[dict] = []
        for file_path in sorted(root.rglob("*")):
            if file_path.is_file() and file_path.suffix.lower() in SUPPORTED_EXTENSIONS:
                try:
                    subfolder = str(file_path.parent.relative_to(root))
                    if subfolder == ".":
                        subfolder = ""
                    rows.append({
                        "File":      file_path.name,
                        "Type":      file_path.suffix.lower().lstrip("."),
                        "Size (KB)": round(file_path.stat().st_size / 1024, 2),
                        "Subfolder": subfolder,
                    })
                except Exception as exc:
                    logger.warning("Could not stat %s: %s", file_path, exc)
        return rows

    def load_single_file(
        self,
        path: Path,
        source_type: str = "local",
    ) -> Optional[RawDocument]:
        return self._extract(path, source_type=source_type)

    def _extract(
        self,
        path: Path,
        source_type: str = "local",
        extra_metadata: Optional[dict] = None,
    ) -> Optional[RawDocument]:
        from utils.helpers import make_doc_id

        ext       = path.suffix.lower()
        extractor = _EXTRACTOR_MAP.get(ext)

        if extractor is None:
            logger.warning("No extractor for file type '%s': %s", ext, path.name)
            return None

        try:
            pages = extractor(path)
            if not pages:
                logger.warning("No text extracted from: %s", path.name)
                return None

            return RawDocument(
                file_path=path,
                doc_id=make_doc_id(path),
                pages=pages,
                total_pages=len(pages),
                source_type=source_type,
                extra_metadata=extra_metadata or {},
            )
        except Exception as exc:
            logger.exception("Failed to process '%s': %s", path.name, exc)
            return None
